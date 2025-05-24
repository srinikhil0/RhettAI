from typing import Dict, List, Optional
import os
from pptx import Presentation
from PyPDF2 import PdfReader
import json
from datetime import datetime

class ContentProcessor:
    def __init__(self, storage_dir: str = "processed_content"):
        self.storage_dir = storage_dir
        os.makedirs(storage_dir, exist_ok=True)
        self.metadata_file = os.path.join(storage_dir, "metadata.json")
        self.load_metadata()

    def load_metadata(self):
        """Load the metadata of processed files."""
        if os.path.exists(self.metadata_file):
            with open(self.metadata_file, 'r') as f:
                self.metadata = json.load(f)
        else:
            self.metadata = {}

    def save_metadata(self):
        """Save the metadata of processed files."""
        with open(self.metadata_file, 'w') as f:
            json.dump(self.metadata, f, indent=2)

    def process_pptx(self, file_path: str) -> Dict:
        """Extract text from PowerPoint files."""
        print(f"Processing PowerPoint file: {file_path}")
        prs = Presentation(file_path)
        content = []
        
        # Extract text from slides
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            content.append({
                "slide_number": i + 1,
                "text": "\n".join(slide_text)
            })
        
        return {
            "type": "presentation",
            "total_slides": len(prs.slides),
            "content": content
        }

    def process_pdf(self, file_path: str) -> Dict:
        """Extract text from PDF files."""
        print(f"Processing PDF file: {file_path}")
        reader = PdfReader(file_path)
        content = []
        
        # Extract text from pages
        for i, page in enumerate(reader.pages):
            content.append({
                "page_number": i + 1,
                "text": page.extract_text()
            })
        
        return {
            "type": "document",
            "total_pages": len(reader.pages),
            "content": content
        }

    def process_file(self, file_path: str, file_metadata: Dict) -> Optional[Dict]:
        """Process a file and store its content."""
        file_name = os.path.basename(file_path)
        file_id = file_metadata['id']
        
        # Check if file is already processed and up to date
        if file_id in self.metadata:
            stored_modified_time = self.metadata[file_id]['modified_time']
            current_modified_time = file_metadata['modifiedTime']
            if stored_modified_time == current_modified_time:
                print(f"File {file_name} is already up to date")
                return None

        # Process based on file type
        if file_path.endswith('.pptx'):
            content = self.process_pptx(file_path)
        elif file_path.endswith('.pdf'):
            content = self.process_pdf(file_path)
        else:
            print(f"Unsupported file type: {file_path}")
            return None

        # Store processed content
        processed_file = {
            "file_name": file_name,
            "file_id": file_id,
            "modified_time": file_metadata['modifiedTime'],
            "processed_time": datetime.now().isoformat(),
            "content": content
        }

        # Save to file
        output_file = os.path.join(self.storage_dir, f"{file_id}.json")
        with open(output_file, 'w') as f:
            json.dump(processed_file, f, indent=2)

        # Update metadata
        self.metadata[file_id] = {
            "file_name": file_name,
            "modified_time": file_metadata['modifiedTime'],
            "processed_time": processed_file['processed_time']
        }
        self.save_metadata()

        print(f"Successfully processed and stored content from {file_name}")
        return processed_file

    def get_processed_content(self, file_id: str) -> Optional[Dict]:
        """Retrieve processed content for a file."""
        if file_id not in self.metadata:
            return None
        
        output_file = os.path.join(self.storage_dir, f"{file_id}.json")
        if not os.path.exists(output_file):
            return None
            
        with open(output_file, 'r') as f:
            return json.load(f) 