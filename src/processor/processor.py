import os
import json
from datetime import datetime
from typing import Dict, List, Optional
from src.database.content_db import ContentDatabase
from pptx import Presentation
from docx import Document
import PyPDF2

class ContentProcessor:
    def __init__(self):
        """Initialize the content processor."""
        self.db = ContentDatabase()
        self.supported_types = {
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._process_presentation,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._process_document,
            'application/pdf': self._process_pdf,
            # Add more MIME types and handlers as needed
        }

    def process_file(self, file_path: str, file_id: str, mime_type: str, modified_time: str) -> bool:
        """Process a file and store its content in the database."""
        try:
            # Check if file type is supported
            if mime_type not in self.supported_types:
                print(f"Unsupported file type: {mime_type}")
                return False

            # Get the appropriate processor for the file type
            processor = self.supported_types[mime_type]
            
            # Process the file
            content = processor(file_path)
            if not content:
                return False

            # Prepare file data for database storage
            file_data = {
                'file_id': file_id,
                'file_name': os.path.basename(file_path),
                'modified_time': modified_time,
                'processed_time': datetime.now().isoformat(),
                'content': content,
                'file_path': file_path
            }

            # Store in database
            self.db.store_file(file_data)
            print(f"Successfully processed and stored: {file_data['file_name']}")
            return True

        except Exception as e:
            print(f"Error processing file {file_path}: {e}")
            return False

    def _process_presentation(self, file_path: str) -> Optional[Dict]:
        """Extract text from PPTX presentations."""
        try:
            prs = Presentation(file_path)
            slides_content = []
            for i, slide in enumerate(prs.slides):
                text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
                slides_content.append({
                    'slide_number': i + 1,
                    'text': "\n".join(text)
                })
            return {
                'type': 'presentation',
                'total_slides': len(prs.slides),
                'content': slides_content
            }
        except Exception as e:
            print(f"Error processing presentation {file_path}: {e}")
            return None

    def _process_document(self, file_path: str) -> Optional[Dict]:
        """Extract text from DOCX documents."""
        try:
            doc = Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return {
                'type': 'document',
                'total_pages': 1,  # DOCX doesn't have pages, treat as 1
                'content': [
                    {
                        'page_number': 1,
                        'text': "\n".join(paragraphs)
                    }
                ]
            }
        except Exception as e:
            print(f"Error processing document {file_path}: {e}")
            return None

    def _process_pdf(self, file_path: str) -> Optional[Dict]:
        """Extract text from PDF files."""
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                content = []
                for i, page in enumerate(reader.pages):
                    text = page.extract_text()
                    content.append({
                        'page_number': i + 1,
                        'text': text or ''
                    })
            return {
                'type': 'document',
                'total_pages': len(content),
                'content': content
            }
        except Exception as e:
            print(f"Error processing PDF {file_path}: {e}")
            return None 